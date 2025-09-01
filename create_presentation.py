#!/usr/bin/env python3
"""
PowerPoint Presentation Generator for Breast Cancer Detection System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Breast Cancer Detection System"
    subtitle.text = "Deep Learning Implementation for Histopathological Image Analysis\n\n[Your Name]\n[University]\n[Date]"
    
    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Problem Statement"
    content.text = """The Challenge:
• Breast cancer: Leading cause of cancer death in women
• Manual diagnosis: Time-consuming, subjective, prone to error
• Pathologist shortage: Limited access to expert diagnosis
• Consistency issues: Inter-observer variability

Our Solution:
AI-powered automated detection system for histopathological images"""
    
    # Slide 3: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    content.text = """What We Built:
• Multi-dataset deep learning system
• 8 cancer subtype classification
• Real-time inference API
• Explainable AI integration

Key Innovation:
First system to combine BreakHis + BACH datasets with explainable AI"""
    
    # Slide 4: Dataset Integration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Dataset Integration"
    content.text = """Multi-Dataset Approach:
• BreakHis: ~7,900 images, 8 subtypes, 40X-400X magnifications
• BACH: 400 images, 4 categories, high-resolution
• Combined: ~8,300 images, unified 5 classes, multi-scale

Why This Matters:
• Improved generalization across imaging conditions
• Larger training dataset for better performance
• Cross-dataset validation for robustness"""
    
    # Slide 5: Technical Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Architecture"
    content.text = """Model Pipeline:
Input Image → Preprocessing → EfficientNetB0 → Classification
   224x224      Normalization    Transfer       8 Classes
    RGB         Augmentation     Learning      + Confidence

Key Components:
• EfficientNetB0: State-of-the-art CNN architecture
• Transfer Learning: Pre-trained on ImageNet
• Patient-wise Splitting: Prevents data leakage
• Weighted Sampling: Handles class imbalance"""
    
    # Slide 6: Accomplishments
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What I've Accomplished ✅"
    content.text = """Core Implementation:
✅ Multi-Dataset Integration - Unified BreakHis + BACH
✅ Robust Data Pipeline - Patient-wise stratified splits
✅ Advanced Model Architecture - EfficientNetB0 with transfer learning
✅ Class Imbalance Solutions - Weighted sampling & loss functions
✅ Production-Ready API - FastAPI with real-time inference
✅ Explainable AI - Grad-CAM + RAG-based explanations"""
    
    # Slide 7: Data Pipeline Innovation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Data Pipeline Innovation"
    content.text = """Patient-Wise Stratification:
• Prevents data leakage
• No patient appears in both train/test sets
• Realistic performance estimates

Class Balancing:
• Weighted Random Sampling: Balances rare cancer types
• Class Weights: Penalizes misclassification of minorities
• Stratified Splits: Maintains distribution across splits"""
    
    # Slide 8: Explainable AI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Explainable AI Integration"
    content.text = """Visual Explanations - Grad-CAM:
• Heatmaps: Show which regions influenced decision
• Overlay Visualization: Highlight suspicious areas
• Clinical Trust: Essential for medical adoption

Textual Explanations - RAG:
• Context-Aware: Retrieves relevant medical knowledge
• Natural Language: Explains predictions in clinical terms
• Confidence Assessment: Provides uncertainty quantification"""
    
    # Slide 9: API Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Production-Ready API"
    content.text = """Features:
• <2 second inference time
• Multi-format image support
• Comprehensive error handling
• RESTful design

Endpoints:
• /api/predict - Main classification
• /api/explain - Detailed analysis
• /api/classes - Available categories

Real-time capabilities with FastAPI backend"""
    
    # Slide 10: System Capabilities
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Capabilities"
    content.text = """Classification Performance:
• 8 Cancer Subtypes: Benign (4) + Malignant (4)
• Multi-Magnification: Handles 40X to 400X
• Real-Time: <2 seconds per image
• Confidence Scoring: Uncertainty quantification

Clinical Features:
• Risk Assessment: High/Low risk categorization
• Top-3 Predictions: Alternative diagnoses
• Visual Attention: Grad-CAM heatmaps
• Textual Explanations: Clinical context"""
    
    # Slide 11: Clinical Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Clinical Impact"
    content.text = """For Healthcare Providers:
• Diagnostic Support: Assists pathologists in decision-making
• Consistency: Reduces inter-observer variability
• Efficiency: Accelerates diagnostic workflow
• Accessibility: Expert-level analysis in resource-limited settings

For Patients:
• Faster Diagnosis: Reduced waiting times
• Improved Accuracy: AI-assisted detection
• Second Opinion: Additional diagnostic confidence
• Early Detection: Better treatment outcomes"""
    
    # Slide 12: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Enhancements"
    content.text = """Ready for Integration:
• GAN-based Data Augmentation - Synthetic sample generation
• Vision Transformer (ViT) - Attention-based architecture
• Multimodal Learning - Clinical metadata integration
• Magnification Robustness - Cross-scale generalization
• RAG Enhancement - Advanced knowledge retrieval

Research Opportunities:
• Federated Learning: Multi-hospital collaboration
• Active Learning: Efficient annotation strategies
• Uncertainty Quantification: Improved confidence estimation"""
    
    # Slide 13: Demo Capabilities
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Demo Capabilities"
    content.text = """Live Demonstration:
1. Image Upload: Drag & drop interface
2. Real-Time Prediction: Instant results
3. Visual Explanation: Grad-CAM heatmaps
4. Confidence Assessment: Uncertainty quantification
5. Risk Categorization: Clinical decision support

System Performance:
• <2 second inference
• Multi-format support
• Batch processing capability
• Concurrent user support"""
    
    # Slide 14: Contributions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contributions & Impact"
    content.text = """Technical Contributions:
• Multi-dataset integration methodology
• Patient-wise validation framework
• Explainable AI for medical imaging
• Production-ready deployment architecture

Research Impact:
• Reproducible framework for future research
• Open-source implementation for community
• Benchmark results on combined datasets
• Clinical applicability demonstration"""
    
    # Slide 15: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    content.text = """What We Achieved:
✅ Complete end-to-end system for breast cancer detection
✅ Multi-dataset training with improved generalization
✅ Explainable AI integration for clinical trust
✅ Production-ready API with real-time inference
✅ Comprehensive evaluation framework

Impact:
First system combining BreakHis + BACH datasets with 
explainable AI for clinical breast cancer detection"""
    
    # Slide 16: Questions
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Questions & Discussion"
    subtitle.text = """Thank You!

Available for Questions:
• Technical implementation details
• Clinical applications
• Future research directions
• Collaboration opportunities

Contact: [your.email@university.edu]"""
    
    # Save presentation
    prs.save('/Users/mac/Documents/GitHub/Thesis_Cancer_Detection/breast_cancer_detection_presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📁 Saved as: breast_cancer_detection_presentation.pptx")

if __name__ == "__main__":
    create_presentation()